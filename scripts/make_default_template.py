#!/usr/bin/env python3
"""
Create templates/default.pptx and templates/default.manifest.json.

Design language derived from ./examples:

- Pure white background on master and all layouts
- Ultra-large bold black title (64pt content, 60pt cover/section)
- Decorative concentric arcs on slide master (upper-right corner, gray, subtle)
- Minimal open layout — no accent bars or filled panels on basic slides
- Light card fills only for comparison and two-column layouts

Runtime relies on startup auto-mapping; this script only prepares the
template asset and a companion manifest artifact for inspection/debugging.

Usage:
    python scripts/make_default_template.py
"""

from __future__ import annotations

import json
import sys
from io import BytesIO
from pathlib import Path

ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT / "src"))

from copy import deepcopy

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.shapes.autoshape import AutoShapeType
from pptx.util import Inches, Pt

_EMU = 914400  # English Metric Units per inch

from pptx_yaml_engine.mapper.service import finalize_manifest, inspect_template, propose_mapping

TEMPLATES_DIR = ROOT / "templates"
PPTX_OUT = TEMPLATES_DIR / "default.pptx"
JSON_OUT = TEMPLATES_DIR / "default.manifest.json"

# Monochromatic editorial palette
WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # main background
INK = RGBColor(0x11, 0x11, 0x11)    # title text – near black
BODY = RGBColor(0x44, 0x44, 0x44)   # body / description text
LINE = RGBColor(0xCC, 0xCC, 0xCC)   # separator lines, card borders
CARD = RGBColor(0xF2, 0xF2, 0xF2)   # card background (comparison, two-col)


def _solid_fill(fill_obj, rgb: RGBColor) -> None:
    fill_obj.solid()
    fill_obj.fore_color.rgb = rgb


def _hide_line(line_obj) -> None:
    line_obj.fill.background()


def _set_defRPr_style(
    tx_body,
    *,
    rgb: RGBColor | None = None,
    size_pt: int | None = None,
    bold: bool | None = None,
) -> None:
    """Inject default run properties into a placeholder's lstStyle."""
    lst_style = tx_body.find(qn("a:lstStyle"))
    if lst_style is None:
        lst_style = etree.SubElement(tx_body, qn("a:lstStyle"))
    lvl1 = lst_style.find(qn("a:lvl1pPr"))
    if lvl1 is None:
        lvl1 = etree.SubElement(lst_style, qn("a:lvl1pPr"))
    defRPr = lvl1.find(qn("a:defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(lvl1, qn("a:defRPr"))

    if rgb is not None:
        for fill in defRPr.findall(qn("a:solidFill")):
            defRPr.remove(fill)
        solid_fill = etree.SubElement(defRPr, qn("a:solidFill"))
        srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
        srgb.set("val", str(rgb))

    if size_pt is not None:
        defRPr.set("sz", str(size_pt * 100))
    if bold is not None:
        defRPr.set("b", "1" if bold else "0")


def _add_object_placeholder(layout, idx: int, *, left: float, top: float, width: float, height: float) -> None:
    """Append a generic OBJECT (content) placeholder to a layout's spTree via raw XML.

    This is required when a stock layout has fewer content placeholders than the
    semantic mapping needs (e.g. Two Content only ships with 2, but three_cards needs 3).
    The element is appended at the end of the spTree; call _configure_placeholder
    immediately after to set typography and fill.
    """
    spTree = layout.shapes._spTree
    ids = [int(el.get("id", "0")) for el in spTree.iter() if el.get("id") and el.get("id").isdigit()]
    shape_id = max(ids, default=100) + 1
    xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f"<p:nvSpPr>"
        f'<p:cNvPr id="{shape_id}" name="Content Placeholder {idx}"/>'
        f"<p:cNvSpPr><a:spLocks noGrp=\"1\"/></p:cNvSpPr>"
        f"<p:nvPr><p:ph idx=\"{idx}\"/></p:nvPr>"
        f"</p:nvSpPr>"
        f"<p:spPr>"
        f"<a:xfrm>"
        f'<a:off x="{int(left * _EMU)}" y="{int(top * _EMU)}"/>'
        f'<a:ext cx="{int(width * _EMU)}" cy="{int(height * _EMU)}"/>'
        f"</a:xfrm>"
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f"</p:spPr>"
        f"<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>"
        f"</p:sp>"
    )
    spTree.append(etree.fromstring(xml))


def _layout_placeholder(layout, ph_idx: int):
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            return ph
    raise KeyError(f"placeholder idx={ph_idx} not found on layout {layout.name!r}")


def _configure_placeholder(
    layout,
    ph_idx: int,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text_rgb: RGBColor | None = None,
    size_pt: int | None = None,
    bold: bool | None = None,
    margins: tuple[float, float, float, float] = (0.0, 0.0, 0.0, 0.0),
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    fill_rgb: RGBColor | None = None,
    line_rgb: RGBColor | None = None,
) -> None:
    ph = _layout_placeholder(layout, ph_idx)
    ph.left = Inches(left)
    ph.top = Inches(top)
    ph.width = Inches(width)
    ph.height = Inches(height)

    if fill_rgb is None:
        ph.fill.background()
    else:
        _solid_fill(ph.fill, fill_rgb)

    if line_rgb is None:
        _hide_line(ph.line)
    else:
        ph.line.color.rgb = line_rgb

    if not hasattr(ph, "text_frame"):
        return

    frame = ph.text_frame
    frame.vertical_anchor = vertical_anchor
    frame.word_wrap = True
    frame.margin_left = Inches(margins[0])
    frame.margin_top = Inches(margins[1])
    frame.margin_right = Inches(margins[2])
    frame.margin_bottom = Inches(margins[3])
    frame.paragraphs[0].alignment = alignment
    _set_defRPr_style(frame._txBody, rgb=text_rgb, size_pt=size_pt, bold=bold)


def _add_shape(
    container,
    shape_type: MSO_AUTO_SHAPE_TYPE,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: RGBColor | None = None,
    line_rgb: RGBColor | None = None,
    line_width_pt: float | None = None,
    send_to_back: bool = True,
) -> None:
    """Add an autoshape to a layout or master via the internal spTree API."""
    shapes = container.shapes
    autoshape = AutoShapeType(shape_type)
    shape_id = shapes._next_shape_id
    sp = shapes._spTree.add_autoshape(
        shape_id,
        f"{autoshape.basename} {shape_id}",
        autoshape.prst,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    if send_to_back:
        shapes._spTree.remove(sp)
        shapes._spTree.insert(2, sp)
    shape = shapes._shape_factory(sp)

    if fill_rgb is None:
        shape.fill.background()
    else:
        _solid_fill(shape.fill, fill_rgb)

    if line_rgb is None:
        _hide_line(shape.line)
    else:
        shape.line.color.rgb = line_rgb
        if line_width_pt is not None:
            shape.line.width = Pt(line_width_pt)


def _add_master_arcs(prs: Presentation) -> None:
    """Add decorative concentric oval arcs to the slide master.

    Arcs are centered at the upper-right corner of the slide (13.333", 0"),
    creating a subtle gray motif in the right portion of every slide.
    Inner arcs are darker; outer arcs fade lighter.
    """
    # (radius_inches, gray_channel_byte, line_width_pt)
    arc_specs = [
        (4.0, 0xBB, 1.5),
        (5.2, 0xC5, 1.25),
        (6.4, 0xCE, 1.0),
        (7.6, 0xD7, 0.75),
        (8.8, 0xDF, 0.75),
        (9.5, 0xE7, 0.5),
    ]
    cx = prs.slide_width.inches   # upper-right corner x (= slide width)
    cy = 0.0                       # upper-right corner y (= slide top)

    for r, gray, lw in arc_specs:
        _add_shape(
            prs.slide_master,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            left=cx - r,
            top=cy - r,
            width=r * 2,
            height=r * 2,
            fill_rgb=None,
            line_rgb=RGBColor(gray, gray, gray),
            line_width_pt=lw,
            send_to_back=True,
        )


# ---------------------------------------------------------------------------
# Layout styling functions
# ---------------------------------------------------------------------------

def _style_cover_layout(layout) -> None:
    """Title Slide – large bold title with subtitle on white."""
    _solid_fill(layout.background.fill, WHITE)
    # Title (CENTER_TITLE, idx=0) – prominent, upper area
    _configure_placeholder(
        layout, 0,
        left=0.55, top=1.8, width=9.8, height=2.6,
        text_rgb=INK, size_pt=60, bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    # Subtitle (SUBTITLE, idx=1)
    _configure_placeholder(
        layout, 1,
        left=0.55, top=4.6, width=8.5, height=1.6,
        text_rgb=BODY, size_pt=18,
    )


def _style_content_layout(layout) -> None:
    """Title and Content – large bold title + open content area."""
    _solid_fill(layout.background.fill, WHITE)
    # Title (TITLE, idx=0) – reduced to 48pt so even 13-char JP titles fit on one line
    _configure_placeholder(
        layout, 0,
        left=0.55, top=0.4, width=11.0, height=1.65,
        text_rgb=INK, size_pt=48, bold=True,
    )
    # Content (OBJECT, idx=1) – starts with comfortable gap below title
    _configure_placeholder(
        layout, 1,
        left=0.55, top=2.2, width=12.3, height=5.05,
        text_rgb=BODY, size_pt=16,
        margins=(0.0, 0.04, 0.0, 0.0),
    )


def _style_section_layout(layout) -> None:
    """Section Header – large centered title for section breaks.

    Default layout[2] positions TITLE at bottom (idx=0 top≈4.82") and
    BODY above it (idx=1 top≈3.18"). We reposition both intentionally.
    """
    _solid_fill(layout.background.fill, WHITE)
    # Body/subtitle above title (idx=1 in default Section Header)
    _configure_placeholder(
        layout, 1,
        left=0.55, top=4.6, width=9.0, height=1.6,
        text_rgb=BODY, size_pt=18,
    )
    # Title below body in template coordinates (prominent, upper area)
    _configure_placeholder(
        layout, 0,
        left=0.55, top=2.0, width=10.0, height=2.3,
        text_rgb=INK, size_pt=60, bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def _style_two_content_layout(layout) -> None:
    """Two Content – redesigned as a true 3-column card layout.

    The stock Two Content layout ships with only 2 OBJECT placeholders, but both
    three_cards_horizontal and three_cards_vertical map here.  We inject a 3rd
    OBJECT placeholder (idx=3) so the mapper assigns one card per column.

    Column geometry (13.333" slide, 0.55" margins, 0.3" gaps):
        card_width = (13.333 - 1.10 - 0.60) / 3 ≈ 3.878"
        col1 left = 0.55   col2 left = 4.73   col3 left = 8.91
    """
    _solid_fill(layout.background.fill, WHITE)

    # Title (TITLE, idx=0)
    _configure_placeholder(
        layout, 0,
        left=0.55, top=0.35, width=11.0, height=1.45,
        text_rgb=INK, size_pt=48, bold=True,
    )

    # Column geometry
    col_w = 3.878
    col_tops = (2.0, 2.0, 2.0)
    col_lefts = (0.55, 4.73, 8.91)
    col_h = 5.2

    # Add the 3rd OBJECT placeholder before configuring (idx 1 & 2 already exist)
    _add_object_placeholder(
        layout, 3,
        left=col_lefts[2], top=col_tops[2], width=col_w, height=col_h,
    )

    # Card background shapes (rounded rectangles, behind placeholders)
    for col_left in col_lefts:
        _add_shape(
            layout, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left=col_left - 0.08, top=col_tops[0] - 0.1, width=col_w + 0.16, height=col_h + 0.2,
            fill_rgb=CARD, line_rgb=LINE, send_to_back=True,
        )

    # Left content (OBJECT, idx=1)
    _configure_placeholder(
        layout, 1,
        left=col_lefts[0], top=col_tops[0], width=col_w, height=col_h,
        text_rgb=BODY, size_pt=14,
        margins=(0.08, 0.06, 0.08, 0.06),
    )
    # Center content (OBJECT, idx=2)
    _configure_placeholder(
        layout, 2,
        left=col_lefts[1], top=col_tops[1], width=col_w, height=col_h,
        text_rgb=BODY, size_pt=14,
        margins=(0.08, 0.06, 0.08, 0.06),
    )
    # Right content (OBJECT, idx=3)
    _configure_placeholder(
        layout, 3,
        left=col_lefts[2], top=col_tops[2], width=col_w, height=col_h,
        text_rgb=BODY, size_pt=14,
        margins=(0.08, 0.06, 0.08, 0.06),
    )


def _style_comparison_layout(layout) -> None:
    """Comparison – title + two card columns with headers and body."""
    _solid_fill(layout.background.fill, WHITE)
    # Title (TITLE, idx=0)
    _configure_placeholder(
        layout, 0,
        left=0.55, top=0.35, width=11.0, height=1.55,
        text_rgb=INK, size_pt=54, bold=True,
    )
    # Left card
    _add_shape(
        layout, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left=0.55, top=2.05, width=5.95, height=5.15,
        fill_rgb=CARD, line_rgb=LINE, send_to_back=True,
    )
    # Right card
    _add_shape(
        layout, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left=6.83, top=2.05, width=5.95, height=5.15,
        fill_rgb=CARD, line_rgb=LINE, send_to_back=True,
    )
    # Left header (BODY, idx=1)
    _configure_placeholder(
        layout, 1,
        left=0.88, top=2.25, width=5.3, height=0.55,
        text_rgb=INK, size_pt=14, bold=True,
        margins=(0.08, 0.04, 0.08, 0.0),
    )
    # Left body (OBJECT, idx=2)
    _configure_placeholder(
        layout, 2,
        left=0.88, top=2.9, width=5.3, height=4.1,
        text_rgb=BODY, size_pt=14,
        margins=(0.08, 0.04, 0.08, 0.04),
    )
    # Right header (BODY, idx=3)
    _configure_placeholder(
        layout, 3,
        left=7.16, top=2.25, width=5.3, height=0.55,
        text_rgb=INK, size_pt=14, bold=True,
        margins=(0.08, 0.04, 0.08, 0.0),
    )
    # Right body (OBJECT, idx=4)
    _configure_placeholder(
        layout, 4,
        left=7.16, top=2.9, width=5.3, height=4.1,
        text_rgb=BODY, size_pt=14,
        margins=(0.08, 0.04, 0.08, 0.04),
    )


def _style_picture_layout(layout) -> None:
    """Picture with Caption – title + large picture area + side caption.

    Default layout[8] positions TITLE at bottom (idx=0) and PICTURE at top
    (idx=1). We reposition all three placeholders for a horizontal split.
    """
    _solid_fill(layout.background.fill, WHITE)
    # Title (TITLE, idx=0) – moved to top
    _configure_placeholder(
        layout, 0,
        left=0.55, top=0.35, width=11.0, height=1.55,
        text_rgb=INK, size_pt=54, bold=True,
    )
    # Picture placeholder (PICTURE, idx=1) – left area
    _configure_placeholder(
        layout, 1,
        left=0.55, top=2.1, width=8.5, height=5.1,
        fill_rgb=CARD, line_rgb=LINE,
    )
    # Caption (BODY, idx=2) – right column
    _configure_placeholder(
        layout, 2,
        left=9.33, top=2.1, width=3.45, height=5.1,
        text_rgb=BODY, size_pt=14,
        margins=(0.0, 0.04, 0.0, 0.0),
    )


def _style_remaining_layouts(prs: Presentation) -> None:
    """Apply white background to remaining stock layouts (5, 6, 7, 9, 10)."""
    for idx in (5, 6, 7, 9, 10):
        try:
            _solid_fill(prs.slide_layouts[idx].background.fill, WHITE)
        except IndexError:
            pass


# ---------------------------------------------------------------------------
# Public build functions
# ---------------------------------------------------------------------------

def build_template() -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # Master: pure white + subtle arc decorations
    _solid_fill(prs.slide_master.background.fill, WHITE)
    _add_master_arcs(prs)

    _style_cover_layout(prs.slide_layouts[0])
    _style_content_layout(prs.slide_layouts[1])
    _style_section_layout(prs.slide_layouts[2])
    _style_two_content_layout(prs.slide_layouts[3])
    _style_comparison_layout(prs.slide_layouts[4])
    _style_picture_layout(prs.slide_layouts[8])
    _style_remaining_layouts(prs)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def build_manifest(template_bytes: bytes) -> dict:
    inspection = inspect_template(template_bytes)
    proposal = propose_mapping(inspection)
    return finalize_manifest(inspection, proposal)


def main() -> None:
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

    template_bytes = build_template()
    PPTX_OUT.write_bytes(template_bytes)

    manifest = build_manifest(template_bytes)
    JSON_OUT.write_text(json.dumps(manifest, ensure_ascii=False, indent=2))

    print(f"Saved {PPTX_OUT.relative_to(ROOT)}")
    print(f"Saved {JSON_OUT.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
