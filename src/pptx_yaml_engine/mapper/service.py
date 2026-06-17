from __future__ import annotations

from collections import defaultdict
from io import BytesIO
from typing import Any, cast

from pptx import Presentation

from pptx_yaml_engine.errors import DomainError, error_dict
from pptx_yaml_engine.layouts import (
    LAYOUT_ALIASES,
    LAYOUT_SPECS,
    SLOT_KIND_PLACEHOLDER_TYPES,
    get_slot_spec,
    slot_ai_placeholder_names,
)
from pptx_yaml_engine.utils.fingerprint import template_fingerprint
from pptx_yaml_engine.utils.layout_names import (
    canonical_builtin_layout_name,
    layout_names_match,
)
from pptx_yaml_engine.utils.pptx import placeholder_type_name
from pptx_yaml_engine.utils.template_bytes import normalize_template_for_python_pptx


def inspect_template(template_bytes: bytes) -> dict[str, Any]:
    try:
        prs = Presentation(BytesIO(normalize_template_for_python_pptx(template_bytes)))
    except Exception as exc:
        raise DomainError("TEMPLATE_OPEN_FAILED", "Unable to open template as a PowerPoint file") from exc

    if len(prs.slide_layouts) == 0:
        raise DomainError("NO_LAYOUTS_FOUND", "Template has no slide layouts")

    slide_width = int(cast(Any, prs.slide_width))
    slide_height = int(cast(Any, prs.slide_height))
    layouts: list[dict[str, Any]] = []

    for layout_index, layout in enumerate(prs.slide_layouts):
        shapes: list[dict[str, Any]] = []
        for shape in layout.shapes:
            left = int(getattr(shape, "left", 0) or 0)
            top = int(getattr(shape, "top", 0) or 0)
            width = int(getattr(shape, "width", 0) or 0)
            height = int(getattr(shape, "height", 0) or 0)
            shape_info: dict[str, Any] = {
                "shape_name": getattr(shape, "name", ""),
                "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "norm_left": left / slide_width if slide_width else 0,
                "norm_top": top / slide_height if slide_height else 0,
                "norm_width": width / slide_width if slide_width else 0,
                "norm_height": height / slide_height if slide_height else 0,
                "norm_center_x": (left + width / 2) / slide_width if slide_width else 0,
                "norm_center_y": (top + height / 2) / slide_height if slide_height else 0,
            }
            if shape_info["is_placeholder"]:
                placeholder_format = shape.placeholder_format
                shape_info.update(
                    {
                        "placeholder_idx": int(placeholder_format.idx),
                        "placeholder_type": placeholder_type_name(placeholder_format.type),
                    }
                )
            shapes.append(shape_info)

        layouts.append(
            {
                "layout_name": layout.name,
                "builtin_layout_name": canonical_builtin_layout_name(layout.name),
                "layout_index": layout_index,
                "shapes": shapes,
                "placeholders": [shape for shape in shapes if shape["is_placeholder"]],
            }
        )

    return {
        "version": 1,
        "template_fingerprint": template_fingerprint(template_bytes),
        "slide_width_emu": slide_width,
        "slide_height_emu": slide_height,
        "layouts": layouts,
    }


AI_PLACEHOLDER_PREFIX = "AI_"

_PLACEHOLDER_KIND_COMPATIBILITY: dict[str, frozenset[str]] = {
    "text": frozenset({"TITLE", "CENTER_TITLE", "SUBTITLE", "BODY", "OBJECT"}),
    "list": frozenset({"TITLE", "CENTER_TITLE", "SUBTITLE", "BODY", "OBJECT"}),
    "table": frozenset({"TABLE"}),
    "chart": frozenset({"CHART"}),
    "icon": frozenset({"PICTURE"}),
}


def _placeholder_binding(shape: dict[str, Any], *, kind: str) -> dict[str, Any]:
    return {
        "kind": kind,
        "placeholder": {
            "idx": shape["placeholder_idx"],
            "type": shape.get("placeholder_type", "UNKNOWN"),
            "shape_name": shape.get("shape_name", ""),
        },
    }


def _is_ai_placeholder_name(value: str) -> bool:
    return value.startswith(AI_PLACEHOLDER_PREFIX)


def _compatible_placeholder_types(kind: str) -> frozenset[str]:
    return _PLACEHOLDER_KIND_COMPATIBILITY.get(kind, frozenset())


def _placeholder_is_compatible(kind: str, placeholder_type: str) -> bool:
    return placeholder_type.upper() in _compatible_placeholder_types(kind)


def _layout_name_candidates(
    semantic: str,
    ruleset: dict[str, Any] | None = None,
) -> list[tuple[str, str]]:
    candidates: list[tuple[str, str]] = [(semantic, "semantic_name")]
    seen = {semantic}
    for alias in LAYOUT_ALIASES.get(semantic, ()):
        if alias in seen:
            continue
        seen.add(alias)
        candidates.append((alias, "explicit_name"))
    if ruleset and isinstance(ruleset.get("aliases"), dict):
        for alias in ruleset["aliases"].get(semantic, []):
            alias_text = str(alias)
            if alias_text in seen:
                continue
            seen.add(alias_text)
            candidates.append((alias_text, "ruleset_alias"))
    return candidates


def _resolve_layout_binding(
    inspection: dict[str, Any],
    semantic: str,
    ruleset: dict[str, Any] | None = None,
) -> tuple[dict[str, Any] | None, str, list[dict[str, Any]]]:
    candidates = _layout_name_candidates(semantic, ruleset)
    matches: list[tuple[dict[str, Any], str]] = []
    for layout in inspection.get("layouts", []):
        for candidate, strategy in candidates:
            if layout_names_match(str(layout.get("layout_name", "")), candidate):
                matches.append((layout, strategy))
                break
    if not matches:
        return (
            None,
            "not_found",
            [
                error_dict(
                    "PPT_LAYOUT_NOT_FOUND",
                    "PowerPoint layout not found for semantic layout",
                    {"layout": semantic, "expected_layout_names": [candidate for candidate, _ in candidates]},
                )
            ],
        )
    if len(matches) > 1:
        return (
            None,
            "ambiguous",
            [
                error_dict(
                    "DUPLICATE_LAYOUT_NAME",
                    "Multiple PowerPoint layouts match the semantic layout",
                    {
                        "layout": semantic,
                        "expected_layout_names": [candidate for candidate, _ in candidates],
                        "matched_layout_names": [layout["layout_name"] for layout, _ in matches],
                    },
                )
            ],
        )
    layout, strategy = matches[0]
    return layout, strategy, []


def _ai_shapes_by_name(layout: dict[str, Any]) -> dict[str, list[dict[str, Any]]]:
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for shape in layout.get("shapes", []):
        shape_name = str(shape.get("shape_name", ""))
        if _is_ai_placeholder_name(shape_name):
            grouped[shape_name].append(shape)
    return dict(grouped)


def _allowed_ai_placeholder_names(semantic: str) -> frozenset[str]:
    return frozenset(
        name
        for slot in LAYOUT_SPECS[semantic].slots
        for name in slot_ai_placeholder_names(semantic, slot.path)
    )


def _validate_layout_ai_contract(
    semantic: str,
    layout: dict[str, Any],
) -> tuple[dict[str, list[dict[str, Any]]], list[dict[str, Any]]]:
    issues: list[dict[str, Any]] = []
    allowed_names = _allowed_ai_placeholder_names(semantic)
    valid_shapes_by_name: dict[str, list[dict[str, Any]]] = {}
    for shape_name, matched_shapes in _ai_shapes_by_name(layout).items():
        placeholder_shapes = [shape for shape in matched_shapes if bool(shape.get("is_placeholder"))]
        if len(matched_shapes) > 1:
            issues.append(
                error_dict(
                    "AI_PLACEHOLDER_DUPLICATED",
                    "AI placeholder name is duplicated within the layout",
                    {
                        "layout": semantic,
                        "ppt_layout_name": layout.get("layout_name"),
                        "shape_name": shape_name,
                        "count": len(matched_shapes),
                    },
                )
            )
        if len(placeholder_shapes) != len(matched_shapes):
            issues.append(
                error_dict(
                    "AI_TARGET_NOT_PLACEHOLDER",
                    "AI placeholder target must be a PowerPoint placeholder",
                    {
                        "layout": semantic,
                        "ppt_layout_name": layout.get("layout_name"),
                        "shape_name": shape_name,
                    },
                )
            )
        if shape_name not in allowed_names:
            issues.append(
                error_dict(
                    "AI_PLACEHOLDER_UNKNOWN",
                    "AI placeholder name is not allowed for the semantic layout",
                    {
                        "layout": semantic,
                        "ppt_layout_name": layout.get("layout_name"),
                        "shape_name": shape_name,
                    },
                )
            )
        if len(matched_shapes) == 1 and len(placeholder_shapes) == 1 and shape_name in allowed_names:
            valid_shapes_by_name[shape_name] = placeholder_shapes
    return valid_shapes_by_name, issues


def _resolve_slot_binding(
    semantic: str,
    slot_path: str,
    layout: dict[str, Any],
    shapes_by_name: dict[str, list[dict[str, Any]]],
) -> tuple[dict[str, Any] | None, list[dict[str, Any]]]:
    slot_spec = get_slot_spec(semantic, slot_path)
    if slot_spec is None:
        return None, []

    expected_names = slot_ai_placeholder_names(semantic, slot_path)
    matched_names = [name for name in expected_names if name in shapes_by_name]
    if not matched_names:
        if slot_spec.required:
            return (
                None,
                [
                    error_dict(
                        "AI_PLACEHOLDER_MISSING",
                        "Required AI placeholder is missing",
                        {
                            "layout": semantic,
                            "ppt_layout_name": layout.get("layout_name"),
                            "slot": slot_path,
                            "expected_names": list(expected_names),
                        },
                    )
                ],
            )
        return None, []

    if len(matched_names) > 1:
        return (
            None,
            [
                error_dict(
                    "AI_PLACEHOLDER_DUPLICATED",
                    "Multiple AI placeholder aliases are present for the same slot",
                    {
                        "layout": semantic,
                        "ppt_layout_name": layout.get("layout_name"),
                        "slot": slot_path,
                        "expected_names": list(expected_names),
                        "matched_names": matched_names,
                    },
                )
            ],
        )

    matched_name = matched_names[0]
    matched_shapes = shapes_by_name[matched_name]
    placeholder_shapes = [shape for shape in matched_shapes if bool(shape.get("is_placeholder"))]
    if len(placeholder_shapes) != len(matched_shapes):
        return (
            None,
            [
                error_dict(
                    "AI_TARGET_NOT_PLACEHOLDER",
                    "AI placeholder target must be a PowerPoint placeholder",
                    {
                        "layout": semantic,
                        "ppt_layout_name": layout.get("layout_name"),
                        "slot": slot_path,
                        "shape_name": matched_name,
                    },
                )
            ],
        )

    if len(placeholder_shapes) > 1:
        return (
            None,
            [
                error_dict(
                    "AI_PLACEHOLDER_DUPLICATED",
                    "AI placeholder name is duplicated within the layout",
                    {
                        "layout": semantic,
                        "ppt_layout_name": layout.get("layout_name"),
                        "slot": slot_path,
                        "shape_name": matched_name,
                    },
                )
            ],
        )

    placeholder = placeholder_shapes[0]
    placeholder_type = str(placeholder.get("placeholder_type", "UNKNOWN")).upper()
    if not _placeholder_is_compatible(slot_spec.kind, placeholder_type):
        return (
            None,
            [
                error_dict(
                    "AI_PLACEHOLDER_INCOMPATIBLE",
                    "AI placeholder type is incompatible with the semantic slot",
                    {
                        "layout": semantic,
                        "ppt_layout_name": layout.get("layout_name"),
                        "slot": slot_path,
                        "shape_name": matched_name,
                        "expected_names": list(expected_names),
                        "kind": slot_spec.kind,
                        "placeholder_type": placeholder_type,
                    },
                )
            ],
        )

    return _placeholder_binding(placeholder, kind=slot_spec.kind), []


def generate_manifest(
    template_bytes: bytes,
    ruleset: dict[str, Any] | None = None,
    overrides: dict[str, Any] | None = None,
) -> dict[str, Any]:
    inspection = inspect_template(template_bytes)
    proposal = propose_mapping(inspection, ruleset)
    return finalize_manifest(inspection, proposal, overrides)


def _required_missing(semantic: str, slots: dict[str, Any]) -> list[str]:
    spec = LAYOUT_SPECS[semantic]
    return [slot.path for slot in spec.slots if slot.required and slot.path not in slots]


def _missing_slot_details(semantic: str, missing: list[str]) -> list[dict[str, Any]]:
    details: list[dict[str, Any]] = []
    for slot_path in missing:
        slot_spec = get_slot_spec(semantic, slot_path)
        if slot_spec is None:
            details.append({"slot": slot_path})
            continue
        details.append(
            {
                "slot": slot_path,
                "kind": slot_spec.kind,
                "expected_ai_names": list(slot_ai_placeholder_names(semantic, slot_path)),
                "compatible_placeholder_types": list(SLOT_KIND_PLACEHOLDER_TYPES[slot_spec.kind]),
            }
        )
    return details


def _missing_semantic_layouts(layouts: dict[str, Any]) -> list[str]:
    return sorted(set(LAYOUT_SPECS) - set(layouts))


def propose_mapping(
    inspection: dict[str, Any], ruleset: dict[str, Any] | None = None
) -> dict[str, Any]:
    proposal_layouts: dict[str, Any] = {}

    for semantic in LAYOUT_SPECS:
        matched_layout, strategy, issues = _resolve_layout_binding(inspection, semantic, ruleset)
        slots: dict[str, Any] = {}
        if matched_layout is not None:
            shapes_by_name, ai_issues = _validate_layout_ai_contract(semantic, matched_layout)
            issues.extend(ai_issues)
            for slot in LAYOUT_SPECS[semantic].slots:
                binding, slot_issues = _resolve_slot_binding(
                    semantic,
                    slot.path,
                    matched_layout,
                    shapes_by_name,
                )
                issues.extend(slot_issues)
                if binding is not None:
                    slots[slot.path] = binding
        missing = _required_missing(semantic, slots)
        proposal_layouts[semantic] = {
            "ppt_layout_name": matched_layout["layout_name"] if matched_layout is not None else None,
            "match_confidence": 1.0 if matched_layout is not None and not issues and not missing else 0.0,
            "layout_match": {"strategy": strategy},
            "slots": slots,
            "required_missing": missing,
            "issues": issues,
            "warnings": [issue["message"] for issue in issues],
            "status": "override_required" if issues or missing else "ready",
        }

    return {
        "version": 1,
        "template_fingerprint": inspection.get("template_fingerprint"),
        "layouts": proposal_layouts,
    }


def _merge_override(layout_binding: dict[str, Any], override: dict[str, Any]) -> dict[str, Any]:
    merged = {
        **layout_binding,
        "layout_match": dict(layout_binding.get("layout_match", {})),
        "issues": list(layout_binding.get("issues", [])),
        "slots": {key: dict(value) for key, value in layout_binding.get("slots", {}).items()},
    }
    if "ppt_layout_name" in override:
        merged["ppt_layout_name"] = override["ppt_layout_name"]
        merged["layout_match"] = {"strategy": "override"}
    for slot_path, slot_override in override.get("slots", {}).items():
        if not isinstance(slot_override, dict):
            continue
        existing = merged["slots"].get(slot_path, {"kind": "text", "placeholder": {}})
        placeholder = dict(existing.get("placeholder", {}))
        if "idx" in slot_override:
            placeholder["idx"] = slot_override["idx"]
        if "type" in slot_override:
            placeholder["type"] = slot_override["type"]
        if "shape_name" in slot_override:
            placeholder["shape_name"] = slot_override["shape_name"]
        merged["slots"][slot_path] = {
            "kind": slot_override.get("kind", existing.get("kind", "text")),
            "placeholder": placeholder,
        }
    return merged


def finalize_manifest(
    inspection: dict[str, Any], proposal: dict[str, Any], overrides: dict[str, Any] | None = None
) -> dict[str, Any]:
    proposal_layouts = proposal.get("layouts", {})
    missing_semantics = _missing_semantic_layouts(proposal_layouts)
    if missing_semantics:
        raise DomainError(
            "TEMPLATE_LAYOUT_CONTRACT_MISMATCH",
            "Proposal does not satisfy the required semantic layout contract",
            {
                "expected": sorted(LAYOUT_SPECS),
                "actual": sorted(proposal_layouts),
                "missing": missing_semantics,
            },
        )
    layout_overrides = (overrides or {}).get("layouts", {}) if isinstance(overrides, dict) else {}
    manifest_layouts: dict[str, Any] = {}

    for semantic in LAYOUT_SPECS:
        binding = proposal_layouts[semantic]
        merged = _merge_override(binding, layout_overrides.get(semantic, {}))
        if not merged.get("ppt_layout_name"):
            raise DomainError(
                "PPT_LAYOUT_NOT_FOUND",
                "PowerPoint layout is unresolved for semantic layout",
                {"layout": semantic},
            )
        missing = _required_missing(semantic, merged.get("slots", {}))
        if missing:
            raise DomainError(
                "AI_PLACEHOLDER_MISSING",
                f"Required AI placeholders are unresolved for layout '{semantic}'",
                {
                    "layout": semantic,
                    "ppt_layout_name": merged.get("ppt_layout_name"),
                    "missing": missing,
                    "missing_placeholders": _missing_slot_details(semantic, missing),
                },
            )
        manifest_layouts[semantic] = {
            "ppt_layout_name": merged["ppt_layout_name"],
            "match_confidence": merged.get("match_confidence", 0),
            "layout_match": merged.get("layout_match", {}),
            "slots": merged.get("slots", {}),
        }

    manifest = {
        "manifest_version": 1,
        "template_fingerprint": inspection.get("template_fingerprint"),
        "layouts": manifest_layouts,
    }
    report = validate_manifest_against_inspection(inspection, manifest)
    if not report["valid"]:
        raise DomainError("MANIFEST_VALIDATION_FAILED", "Manifest does not match template inspection", {"issues": report["issues"]})
    return manifest


def validate_manifest_against_inspection(inspection: dict[str, Any], manifest: dict[str, Any]) -> dict[str, Any]:
    issues: list[dict[str, Any]] = []
    manifest_layouts = manifest.get("layouts", {})
    missing_semantics = _missing_semantic_layouts(manifest_layouts)
    if missing_semantics:
        issues.append(
            error_dict(
                "TEMPLATE_LAYOUT_CONTRACT_MISMATCH",
                "Manifest does not satisfy the required semantic layout contract",
                {
                    "expected": sorted(LAYOUT_SPECS),
                    "actual": sorted(manifest_layouts),
                    "missing": missing_semantics,
                },
            )
        )
    for semantic, binding in manifest_layouts.items():
        if semantic not in LAYOUT_SPECS:
            issues.append(
                error_dict(
                    "SEMANTIC_LAYOUT_NOT_FOUND",
                    "Semantic layout is not supported",
                    {"layout": semantic},
                )
            )
            continue
        exact_matches = [layout for layout in inspection.get("layouts", []) if layout.get("layout_name") == binding.get("ppt_layout_name")]
        if exact_matches:
            matches = exact_matches
        else:
            matches = [
                layout
                for layout in inspection.get("layouts", [])
                if layout_names_match(str(layout.get("layout_name", "")), str(binding.get("ppt_layout_name", "")))
            ]
        if not matches:
            issues.append(error_dict("PPT_LAYOUT_NOT_FOUND", "PowerPoint layout not found", {"layout": semantic, "ppt_layout_name": binding.get("ppt_layout_name")}))
            continue
        if len(matches) > 1:
            issues.append(
                error_dict(
                    "DUPLICATE_LAYOUT_NAME",
                    "PowerPoint layout name is duplicated",
                    {"layout": semantic, "ppt_layout_name": binding.get("ppt_layout_name")},
                )
            )
            continue
        layout = matches[0]
        shapes_by_name, ai_issues = _validate_layout_ai_contract(semantic, layout)
        issues.extend(ai_issues)
        slots = binding.get("slots", {})
        missing = _required_missing(semantic, slots)
        if missing:
            issues.append(
                error_dict(
                    "AI_PLACEHOLDER_MISSING",
                    "Manifest is missing required AI placeholder bindings",
                    {
                        "layout": semantic,
                        "ppt_layout_name": binding.get("ppt_layout_name"),
                        "missing": missing,
                    },
                )
            )
        for slot_path, slot in slots.items():
            slot_spec = get_slot_spec(semantic, slot_path)
            if slot_spec is None:
                issues.append(
                    error_dict(
                        "SEMANTIC_LAYOUT_NOT_FOUND",
                        "Semantic slot is not supported",
                        {"layout": semantic, "slot": slot_path},
                    )
                )
                continue
            placeholder = slot.get("placeholder", {})
            shape_name = str(placeholder.get("shape_name", ""))
            idx = placeholder.get("idx")
            expected_names = slot_ai_placeholder_names(semantic, slot_path)
            if not shape_name or shape_name not in expected_names:
                issues.append(
                    error_dict(
                        "AI_PLACEHOLDER_MISSING",
                        "Manifest shape_name does not match the AI placeholder contract",
                        {
                            "layout": semantic,
                            "slot": slot_path,
                            "shape_name": shape_name,
                            "expected_names": list(expected_names),
                        },
                    )
                )
                continue
            matched_shapes = shapes_by_name.get(shape_name, [])
            if not matched_shapes:
                issues.append(
                    error_dict(
                        "AI_PLACEHOLDER_MISSING",
                        "AI placeholder referenced by the manifest is missing",
                        {"layout": semantic, "slot": slot_path, "shape_name": shape_name},
                    )
                )
                continue
            placeholder_shapes = [shape for shape in matched_shapes if bool(shape.get("is_placeholder"))]
            if len(placeholder_shapes) != len(matched_shapes):
                issues.append(
                    error_dict(
                        "AI_TARGET_NOT_PLACEHOLDER",
                        "AI placeholder target must be a PowerPoint placeholder",
                        {"layout": semantic, "slot": slot_path, "shape_name": shape_name},
                    )
                )
                continue
            if len(placeholder_shapes) > 1:
                issues.append(
                    error_dict(
                        "AI_PLACEHOLDER_DUPLICATED",
                        "AI placeholder name is duplicated within the layout",
                        {"layout": semantic, "slot": slot_path, "shape_name": shape_name},
                    )
                )
                continue
            actual_placeholder = placeholder_shapes[0]
            placeholder_type = str(actual_placeholder.get("placeholder_type", "UNKNOWN")).upper()
            kind = str(slot.get("kind", slot_spec.kind))
            if not _placeholder_is_compatible(kind, placeholder_type):
                issues.append(
                    error_dict(
                        "AI_PLACEHOLDER_INCOMPATIBLE",
                        "AI placeholder type is incompatible with the semantic slot",
                        {
                            "layout": semantic,
                            "slot": slot_path,
                            "shape_name": shape_name,
                            "kind": kind,
                            "placeholder_type": placeholder_type,
                        },
                    )
                )
            if idx != actual_placeholder.get("placeholder_idx"):
                issues.append(
                    error_dict(
                        "AI_PLACEHOLDER_IDX_MISMATCH",
                        "Manifest idx does not match the authoritative AI placeholder name",
                        {
                            "layout": semantic,
                            "slot": slot_path,
                            "shape_name": shape_name,
                            "expected_idx": actual_placeholder.get("placeholder_idx"),
                            "actual_idx": idx,
                        },
                    )
                )
    return {"valid": not issues, "issues": issues}


def validate_manifest(template_bytes: bytes, manifest: dict[str, Any]) -> dict[str, Any]:
    inspection = inspect_template(template_bytes)
    issues: list[dict[str, Any]] = []
    if manifest.get("template_fingerprint") != inspection.get("template_fingerprint"):
        issues.append(
            error_dict(
                "TEMPLATE_FINGERPRINT_MISMATCH",
                "Manifest fingerprint does not match template",
                {
                    "expected": manifest.get("template_fingerprint"),
                    "actual": inspection.get("template_fingerprint"),
                },
            )
        )
    issues.extend(validate_manifest_against_inspection(inspection, manifest)["issues"])
    return {"valid": not issues, "issues": issues}


def layout_placeholder_counts(inspection: dict[str, Any]) -> dict[str, dict[str, int]]:
    counts: dict[str, dict[str, int]] = {}
    for layout in inspection.get("layouts", []):
        type_counts: defaultdict[str, int] = defaultdict(int)
        for placeholder in layout.get("placeholders", []):
            type_counts[str(placeholder.get("placeholder_type", "UNKNOWN"))] += 1
        counts[layout["layout_name"]] = dict(type_counts)
    return counts
