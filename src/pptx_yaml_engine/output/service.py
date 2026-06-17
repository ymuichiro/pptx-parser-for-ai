from __future__ import annotations

from io import BytesIO
from typing import Any
from unicodedata import east_asian_width

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

from pptx_yaml_engine.errors import DomainError
from pptx_yaml_engine.icons.registry import resolve_icon
from pptx_yaml_engine.output.validation import validate_deck as validate_deck
from pptx_yaml_engine.utils.fingerprint import template_fingerprint
from pptx_yaml_engine.utils.pptx import clear_existing_slides, find_layout_by_name
from pptx_yaml_engine.utils.template_bytes import normalize_template_for_python_pptx


def _get_path(data: dict[str, Any], path: str) -> Any:
    current: Any = data
    for part in path.replace("]", "").split("."):
        if "[" in part:
            name, index_raw = part.split("[", 1)
            current = current.get(name) if isinstance(current, dict) else None
            if not isinstance(current, list):
                return None
            index = int(index_raw)
            if index >= len(current):
                return None
            current = current[index]
        elif isinstance(current, dict):
            current = current.get(part)
        else:
            return None
    return current


def _list_item_text(item: Any) -> str:
    if isinstance(item, dict):
        text = str(item.get("text", ""))
        level = int(item.get("level", 0) or 0)
        return f"{'  ' * max(0, level)}{text}"
    return str(item)


def _slot_value(slide_spec: dict[str, Any], path: str) -> Any:
    if path.endswith(".combined_text"):
        base = path.rsplit(".", 1)[0]
        combined = _get_path(slide_spec, base)
        if isinstance(combined, dict):
            lines: list[dict[str, Any]] = []
            label = combined.get("label")
            title = combined.get("title")
            description = combined.get("description")
            if label is not None and str(label).strip():
                lines.append({"text": str(label).strip(), "level": 1})
            if title is not None and str(title).strip():
                lines.append({"text": str(title).strip(), "level": 0})
            if description is not None and str(description).strip():
                lines.append({"text": str(description).strip(), "level": 1})
            return lines or None
        return None
    return _get_path(slide_spec, path)


def _placeholder(slide: Any, idx: int, *, layout: str, slot: str) -> Any:
    try:
        return slide.placeholders[idx]
    except KeyError as exc:
        raise DomainError(
            "PLACEHOLDER_IDX_NOT_FOUND",
            "Placeholder idx not found on generated slide",
            {"layout": layout, "slot": slot, "idx": idx},
        ) from exc


def _placeholder_binding_name(placeholder_info: dict[str, Any]) -> str | None:
    name = placeholder_info.get("name") or placeholder_info.get("shape_name")
    if name is None:
        return None
    normalized = str(name).strip()
    return normalized or None


def _rename_placeholder(placeholder: Any, name: str | None) -> None:
    if name is None:
        return
    try:
        placeholder.name = name
    except AttributeError:
        return


def _clear_bullets(paragraph: Any) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag.rsplit("}", 1)[-1].startswith("bu"):
            p_pr.remove(child)
    p_pr.insert(0, OxmlElement("a:buNone"))


FONT_FAMILY = "Aptos"
CJK_FONT_FAMILY = "Yu Gothic"


def _rgb(red: int, green: int, blue: int) -> Any:
    return RGBColor(red, green, blue)  # type: ignore[no-untyped-call]


THEME_BLACK = _rgb(0x19, 0x19, 0x19)
THEME_CHARCOAL = _rgb(0x24, 0x24, 0x24)
THEME_MUTED = _rgb(0x66, 0x66, 0x66)
THEME_LIGHT_GRAY = _rgb(0xF7, 0xF7, 0xF5)
THEME_ROW_GRAY = _rgb(0xEF, 0xEF, 0xEC)
THEME_LINE_GRAY = _rgb(0xD9, 0xD9, 0xD6)
THEME_TEXT = THEME_CHARCOAL
THEME_WHITE = _rgb(0xFF, 0xFF, 0xFF)
CHART_SERIES_COLORS = (
    THEME_BLACK,
    THEME_CHARCOAL,
    _rgb(0x74, 0x74, 0x74),
    _rgb(0x9A, 0x9A, 0x97),
    _rgb(0xBC, 0xBC, 0xB8),
)

TEXT_STYLE_OVERRIDES: dict[tuple[str, str], tuple[float, bool]] = {
    ("cover_title", "title"): (40, True),
    ("cover_title", "subtitle"): (16, False),
    ("cover_title", "date"): (10, False),
    ("cover_title", "organization"): (10.5, False),
    ("cover_title", "author"): (10.5, False),
    ("section_divider", "title"): (34, True),
    ("section_divider", "subtitle"): (14.5, False),
    ("section_divider", "section_no"): (86, True),
    ("agenda", "title"): (29, True),
    ("agenda", "subtitle"): (12, False),
    ("agenda", "items"): (16.5, False),
    ("comparison_2col", "title"): (29, True),
    ("comparison_2col", "subtitle"): (10.5, False),
    ("comparison_2col", "left.title"): (13.5, True),
    ("comparison_2col", "right.title"): (13.5, True),
    ("comparison_2col", "left.description"): (12.5, False),
    ("comparison_2col", "right.description"): (12.5, False),
    ("comparison_2col", "left.bullets"): (11.5, False),
    ("comparison_2col", "right.bullets"): (11.5, False),
    ("three_cards_vertical", "title"): (29, True),
    ("three_cards_vertical", "subtitle"): (10, False),
    ("three_cards_vertical", "cards[0].title"): (13.5, True),
    ("three_cards_vertical", "cards[1].title"): (13.5, True),
    ("three_cards_vertical", "cards[2].title"): (13.5, True),
    ("three_cards_vertical", "cards[0].description"): (11.5, False),
    ("three_cards_vertical", "cards[1].description"): (11.5, False),
    ("three_cards_vertical", "cards[2].description"): (11.5, False),
    ("chart_basic", "title"): (29, True),
    ("chart_basic", "subtitle"): (11, False),
    ("chart_basic", "caption"): (10, False),
    ("table_basic", "title"): (29, True),
    ("table_basic", "subtitle"): (11, False),
    ("table_basic", "caption"): (10, False),
    ("list_basic", "title"): (29, True),
    ("list_basic", "subtitle"): (11, False),
    ("list_basic", "items"): (15.5, False),
    ("image_caption", "title"): (29, True),
    ("image_caption", "subtitle"): (11.5, False),
    ("image_caption", "caption"): (12, False),
    ("image_caption", "attribution"): (10.5, False),
    ("appendix_backup", "title"): (29, True),
    ("appendix_backup", "subtitle"): (11, False),
    ("appendix_backup", "body"): (12.5, False),
    ("appendix_backup", "items"): (11, False),
    ("appendix_backup", "references"): (12, False),
    ("closing_end", "title"): (32, True),
    ("closing_end", "subtitle"): (14.5, False),
    ("closing_end", "message"): (14, False),
    ("closing_end", "contact"): (11.5, False),
    ("closing_end", "cta"): (11.5, True),
}

TEXT_COLOR_OVERRIDES: dict[tuple[str, str], Any] = {
    ("cover_title", "date"): THEME_MUTED,
    ("cover_title", "organization"): THEME_MUTED,
    ("cover_title", "author"): THEME_MUTED,
    ("section_divider", "subtitle"): THEME_MUTED,
    ("agenda", "subtitle"): THEME_MUTED,
    ("comparison_2col", "subtitle"): THEME_MUTED,
    ("comparison_2col", "left.description"): THEME_MUTED,
    ("comparison_2col", "right.description"): THEME_MUTED,
    ("comparison_2col", "left.bullets"): THEME_MUTED,
    ("comparison_2col", "right.bullets"): THEME_MUTED,
    ("three_cards_vertical", "subtitle"): THEME_MUTED,
    ("three_cards_vertical", "cards[0].description"): THEME_MUTED,
    ("three_cards_vertical", "cards[1].description"): THEME_MUTED,
    ("three_cards_vertical", "cards[2].description"): THEME_MUTED,
    ("chart_basic", "subtitle"): THEME_MUTED,
    ("chart_basic", "caption"): THEME_MUTED,
    ("table_basic", "subtitle"): THEME_MUTED,
    ("table_basic", "caption"): THEME_MUTED,
    ("list_basic", "subtitle"): THEME_MUTED,
    ("image_caption", "subtitle"): THEME_MUTED,
    ("image_caption", "attribution"): THEME_MUTED,
    ("appendix_backup", "subtitle"): THEME_MUTED,
    ("appendix_backup", "items"): THEME_MUTED,
    ("appendix_backup", "references"): THEME_MUTED,
    ("closing_end", "contact"): THEME_MUTED,
}

CJK_MAJOR_TITLE_MIN_SIZES: dict[tuple[str, str], float] = {
    ("cover_title", "title"): 36.0,
    ("section_divider", "title"): 31.0,
    ("agenda", "title"): 27.0,
    ("comparison_2col", "title"): 27.0,
    ("three_cards_vertical", "title"): 27.0,
    ("chart_basic", "title"): 27.0,
    ("table_basic", "title"): 27.0,
    ("list_basic", "title"): 27.0,
    ("image_caption", "title"): 27.0,
    ("appendix_backup", "title"): 27.0,
    ("closing_end", "title"): 29.0,
}

TEXT_ALIGNMENT_OVERRIDES: dict[tuple[str, str], PP_ALIGN] = {
    ("three_cards_vertical", "cards[0].title"): PP_ALIGN.CENTER,
    ("three_cards_vertical", "cards[1].title"): PP_ALIGN.CENTER,
    ("three_cards_vertical", "cards[2].title"): PP_ALIGN.CENTER,
}


def _configure_text_frame(frame: Any) -> None:
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = Pt(2)
    frame.margin_right = Pt(2)
    frame.margin_top = Pt(1)
    frame.margin_bottom = Pt(1)


def _contains_cjk(text: str) -> bool:
    return any(
        "\u3040" <= char <= "\u30ff"
        or "\u3400" <= char <= "\u9fff"
        or "\uf900" <= char <= "\ufaff"
        or "\uff00" <= char <= "\uffef"
        for char in text
    )


def _text_display_width(text: str) -> int:
    return sum(2 if east_asian_width(char) in {"F", "W"} else 1 for char in text)


def _font_family_for_text(text: str) -> str:
    return CJK_FONT_FAMILY if _contains_cjk(text) else FONT_FAMILY


def _apply_font_family(font: Any, family: str) -> None:
    font.name = family
    r_pr = getattr(font, "_element", None)
    if r_pr is None:
        return
    for tag in ("a:latin", "a:ea", "a:cs"):
        child = r_pr.find(qn(tag))
        if child is None:
            child = OxmlElement(tag)
            r_pr.insert(0, child)
        child.set("typeface", family)


def _cjk_major_title_adjustment(layout: str, slot_path: str, max_line_length: int) -> float:
    if (layout, slot_path) not in CJK_MAJOR_TITLE_MIN_SIZES:
        return 0.0
    if max_line_length >= 36:
        return 3.0
    if max_line_length >= 28:
        return 2.0
    if max_line_length >= 24:
        return 1.5
    return 0.0


def _font_size(
    layout: str,
    slot_path: str,
    *,
    kind: str,
    line_count: int,
    max_line_length: int,
    contains_cjk: bool = False,
) -> tuple[float, bool]:
    size, bold = TEXT_STYLE_OVERRIDES.get((layout, slot_path), (11.0 if kind == "list" else 10.0, False))
    if line_count >= 5:
        size -= 0.8
    if max_line_length >= 86:
        size -= 1.0
    elif max_line_length >= 64:
        size -= 0.5
    if contains_cjk:
        size -= _cjk_major_title_adjustment(layout, slot_path, max_line_length)
        size = max(CJK_MAJOR_TITLE_MIN_SIZES.get((layout, slot_path), 5.0), size)
    return max(5.0, size), bold


def _font_color(layout: str, slot_path: str) -> Any:
    return TEXT_COLOR_OVERRIDES.get((layout, slot_path), THEME_TEXT)


def _paragraph_alignment(layout: str, slot_path: str) -> PP_ALIGN | None:
    return TEXT_ALIGNMENT_OVERRIDES.get((layout, slot_path))


def _style_paragraph(paragraph: Any, *, size: float, bold: bool, color: Any, keep_bullets: bool, family: str) -> None:
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(4 if keep_bullets else 0)
    paragraph.line_spacing = 1.08 if keep_bullets else (0.98 if size <= 9 else 1.02)
    _apply_font_family(paragraph.font, family)
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    for run in paragraph.runs:
        _apply_font_family(run.font, family)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _style_table_cell(cell: Any, *, header: bool, row_index: int) -> None:
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = THEME_BLACK if header else (THEME_ROW_GRAY if row_index % 2 else THEME_LIGHT_GRAY)
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)
    family = _font_family_for_text(cell.text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 0.95
        _apply_font_family(paragraph.font, family)
        paragraph.font.size = Pt(11.5 if header else 11)
        paragraph.font.bold = header
        paragraph.font.color.rgb = THEME_WHITE if header else THEME_TEXT
        for run in paragraph.runs:
            _apply_font_family(run.font, family)
            run.font.size = Pt(11.5 if header else 11)
            run.font.bold = header
            run.font.color.rgb = THEME_WHITE if header else THEME_TEXT


def _icon_ref_for_slot(icon_ref: dict[str, Any], *, layout: str, slot_path: str) -> dict[str, Any]:
    if "color" in icon_ref:
        return icon_ref
    styled = dict(icon_ref)
    if (layout == "three_cards_vertical" and slot_path.endswith(".icon")) or layout == "comparison_2col":
        styled["color"] = "#FFFFFF"
    elif layout == "image_caption":
        styled["color"] = "#242424"
    return styled


def _text_lines(value: Any) -> list[tuple[str, int]]:
    if value is None:
        return []
    if isinstance(value, list):
        lines: list[tuple[str, int]] = []
        for item in value:
            if isinstance(item, dict) and "text" in item:
                text = str(item.get("text", "")).strip()
                if text:
                    lines.append((text, int(item.get("level", 0) or 0)))
                continue
            lines.extend(_text_lines(_list_item_text(item)))
        return lines
    if isinstance(value, dict):
        return [(str(item).strip(), 0) for item in value.values() if item is not None and str(item).strip()]
    return [(line.strip(), 0) for line in str(value).splitlines() if line.strip()]


def _clear_text_placeholder(placeholder: Any) -> None:
    if not hasattr(placeholder, "text_frame"):
        return
    placeholder.text_frame.clear()


def _write_text(placeholder: Any, value: Any, *, layout: str, slot_path: str, append: bool = False) -> None:
    if not hasattr(placeholder, "text_frame"):
        return
    lines = _text_lines(value)
    if not lines:
        return

    frame = placeholder.text_frame
    _configure_text_frame(frame)
    if not append:
        frame.clear()
    combined_text = "\n".join(line for line, _level in lines)
    font_size, bold = _font_size(
        layout,
        slot_path,
        kind="text",
        line_count=len(lines),
        max_line_length=max(_text_display_width(line) for line, _level in lines),
        contains_cjk=_contains_cjk(combined_text),
    )
    font_color = _font_color(layout, slot_path)
    first_paragraph = frame.paragraphs[0] if (not append or not frame.text.strip()) else None
    family = _font_family_for_text(combined_text)
    for line, level in lines:
        paragraph = first_paragraph if first_paragraph is not None else frame.add_paragraph()
        first_paragraph = None
        paragraph.text = line
        paragraph.level = level
        paragraph.alignment = _paragraph_alignment(layout, slot_path)
        _clear_bullets(paragraph)
        _style_paragraph(paragraph, size=font_size, bold=bold, color=font_color, keep_bullets=False, family=family)


def _write_list(placeholder: Any, items: list[Any], *, layout: str, slot_path: str, append: bool = False) -> None:
    if not hasattr(placeholder, "text_frame"):
        return
    frame = placeholder.text_frame
    _configure_text_frame(frame)
    if not append:
        frame.clear()
    item_lines = [_list_item_text(item) for item in items]
    font_size, bold = _font_size(
        layout,
        slot_path,
        kind="list",
        line_count=len(item_lines),
        max_line_length=max((_text_display_width(line) for line in item_lines), default=0),
        contains_cjk=_contains_cjk("\n".join(item_lines)),
    )
    font_color = _font_color(layout, slot_path)
    first_paragraph = frame.paragraphs[0] if (not append or not frame.text.strip()) else None
    family = _font_family_for_text("\n".join(item_lines))
    for item in items:
        paragraph = first_paragraph if first_paragraph is not None else frame.add_paragraph()
        first_paragraph = None
        paragraph.text = _list_item_text(item)
        if isinstance(item, dict):
            paragraph.level = int(item.get("level", 0) or 0)
        _style_paragraph(paragraph, size=font_size, bold=bold, color=font_color, keep_bullets=True, family=family)


def _write_icon(slide: Any, placeholder: Any, icon_ref: dict[str, Any], *, layout: str, slot_path: str) -> None:
    target_px = min(512, max(64, int(max(int(placeholder.width), int(placeholder.height)) / 9525)))
    icon_bytes = resolve_icon(_icon_ref_for_slot(icon_ref, layout=layout, slot_path=slot_path), target_px=target_px)
    stream = BytesIO(icon_bytes)
    if hasattr(placeholder, "insert_picture"):
        placeholder.insert_picture(stream)
        return
    raise DomainError(
        "PLACEHOLDER_INSERT_UNSUPPORTED",
        "Icon slot must be a picture placeholder that supports insert_picture()",
        {"placeholder_type": type(placeholder).__name__},
    )


def _write_table(slide: Any, placeholder: Any, table_spec: dict[str, Any]) -> None:
    headers = [str(value) for value in table_spec["headers"]]
    rows = [[str(value) for value in row] for row in table_spec.get("rows", [])]
    row_count = len(rows) + 1
    col_count = len(headers)
    if hasattr(placeholder, "insert_table"):
        frame = placeholder.insert_table(row_count, col_count)
    else:
        raise DomainError(
            "PLACEHOLDER_INSERT_UNSUPPORTED",
            "Table slot must be a table placeholder that supports insert_table()",
            {"placeholder_type": type(placeholder).__name__},
        )
    table = frame.table
    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header
        _style_table_cell(table.cell(0, col_idx), header=True, row_index=0)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = value
            _style_table_cell(table.cell(row_idx, col_idx), header=False, row_index=row_idx)


def _chart_type(kind: str) -> Any:
    return {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }[kind]


def _chart_color(index: int) -> Any:
    return CHART_SERIES_COLORS[index % len(CHART_SERIES_COLORS)]


def _style_chart_line(line: Any, color: Any, *, width_pt: float = 1.0) -> None:
    line.color.rgb = color
    line.width = Pt(width_pt)


def _style_chart_series(chart: Any) -> None:
    for index, series in enumerate(chart.series):
        color = _chart_color(index)
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        _style_chart_line(series.format.line, color, width_pt=1.0)
        if hasattr(series, "marker"):
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = color
            _style_chart_line(series.marker.format.line, color, width_pt=0.75)
        for point_index, point in enumerate(getattr(series, "points", [])):
            point_color = color if len(chart.series) > 1 else _chart_color(point_index)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = point_color
            _style_chart_line(point.format.line, THEME_WHITE, width_pt=0.75)


def _style_chart_axis(axis: Any, *, gridlines: bool) -> None:
    axis.tick_labels.font.name = FONT_FAMILY
    axis.tick_labels.font.size = Pt(9)
    axis.tick_labels.font.color.rgb = THEME_MUTED
    _style_chart_line(axis.format.line, THEME_LINE_GRAY, width_pt=0.6)
    if gridlines and getattr(axis, "has_major_gridlines", False):
        _style_chart_line(axis.major_gridlines.format.line, THEME_LINE_GRAY, width_pt=0.45)


def _try_style_chart_axis(chart: Any, axis_name: str, *, gridlines: bool) -> None:
    try:
        axis = getattr(chart, axis_name)
    except ValueError:
        return
    _style_chart_axis(axis, gridlines=gridlines)


def _style_chart(chart: Any) -> None:
    chart.chart_style = None
    chart.font.name = FONT_FAMILY
    chart.font.size = Pt(9)
    chart.font.color.rgb = THEME_MUTED
    if chart.has_legend:
        chart.legend.font.name = FONT_FAMILY
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = THEME_MUTED
    _style_chart_series(chart)
    _try_style_chart_axis(chart, "category_axis", gridlines=False)
    _try_style_chart_axis(chart, "value_axis", gridlines=True)


def _write_chart(slide: Any, placeholder: Any, chart_spec: dict[str, Any]) -> None:
    chart_data = CategoryChartData()  # type: ignore[no-untyped-call]
    chart_data.categories = [str(value) for value in chart_spec["categories"]]
    for series in chart_spec["series"]:
        chart_data.add_series(str(series["name"]), tuple(series["values"]))  # type: ignore[no-untyped-call]
    chart_type = _chart_type(chart_spec["kind"])
    if hasattr(placeholder, "insert_chart"):
        graphic_frame = placeholder.insert_chart(chart_type, chart_data)
        _style_chart(graphic_frame.chart)
    else:
        raise DomainError(
            "PLACEHOLDER_INSERT_UNSUPPORTED",
            "Chart slot must be a chart placeholder that supports insert_chart()",
            {"placeholder_type": type(placeholder).__name__},
        )


def _render_slide(slide: Any, slide_spec: dict[str, Any], binding: dict[str, Any]) -> None:
    seen_idx: set[int] = set()
    for slot_path, slot_binding in binding.get("slots", {}).items():
        value = _slot_value(slide_spec, slot_path)
        placeholder_info = slot_binding.get("placeholder", {})
        idx = int(placeholder_info["idx"])
        placeholder = _placeholder(slide, idx, layout=slide_spec["layout"], slot=slot_path)
        kind = slot_binding.get("kind", "text")
        if idx not in seen_idx:
            _rename_placeholder(placeholder, _placeholder_binding_name(placeholder_info))
        if value is None:
            if kind in {"text", "list"} and idx not in seen_idx:
                _clear_text_placeholder(placeholder)
                seen_idx.add(idx)
            continue
        append = idx in seen_idx
        seen_idx.add(idx)
        if kind == "list":
            _write_list(placeholder, value if isinstance(value, list) else [value], layout=slide_spec["layout"], slot_path=slot_path, append=append)
        elif kind == "table":
            _write_table(slide, placeholder, value)
        elif kind == "chart":
            _write_chart(slide, placeholder, value)
        elif kind == "icon":
            _write_icon(slide, placeholder, value, layout=slide_spec["layout"], slot_path=slot_path)
        else:
            _write_text(placeholder, value, layout=slide_spec["layout"], slot_path=slot_path, append=append)


def render_pptx(template_bytes: bytes, manifest: dict[str, Any], deck: dict[str, Any]) -> bytes:
    if manifest.get("template_fingerprint") != template_fingerprint(template_bytes):
        raise DomainError(
            "TEMPLATE_FINGERPRINT_MISMATCH",
            "Manifest fingerprint does not match template",
            {
                "expected": manifest.get("template_fingerprint"),
                "actual": template_fingerprint(template_bytes),
            },
        )
    validation = validate_deck(deck, manifest)
    if not validation["valid"]:
        raise DomainError(
            "DECK_SCHEMA_INVALID",
            "Deck validation failed",
            {"issues": validation["issues"]},
        )

    try:
        prs = Presentation(BytesIO(normalize_template_for_python_pptx(template_bytes)))
    except Exception as exc:
        raise DomainError(
            "TEMPLATE_OPEN_FAILED",
            "Unable to open template as a PowerPoint file",
        ) from exc

    clear_existing_slides(prs)
    for slide_spec in deck["slides"]:
        layout_name = manifest["layouts"][slide_spec["layout"]]["ppt_layout_name"]
        layout = find_layout_by_name(prs, layout_name)
        if layout is None:
            raise DomainError(
                "PPT_LAYOUT_NOT_FOUND",
                "PowerPoint layout not found",
                {"layout": slide_spec["layout"], "ppt_layout_name": layout_name},
            )
        if isinstance(layout, list):
            raise DomainError(
                "DUPLICATE_LAYOUT_NAME",
                "PowerPoint layout name is duplicated",
                {"ppt_layout_name": layout_name},
            )
        slide = prs.slides.add_slide(layout)
        binding = manifest["layouts"][slide_spec["layout"]]
        _render_slide(slide, slide_spec, binding)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
