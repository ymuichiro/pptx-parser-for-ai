from __future__ import annotations

import json
from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

import pytest
import yaml  # type: ignore[import-untyped]
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from pptx_yaml_engine.layouts import LAYOUT_SPECS, slot_ai_placeholder_names
from pptx_yaml_engine.mapper.service import generate_manifest, inspect_template, validate_manifest
from pptx_yaml_engine.output.service import (
    CJK_FONT_FAMILY,
    FONT_FAMILY,
    _text_display_width,
    render_pptx,
)
from pptx_yaml_engine.server.template_registry import TemplateRegistry
from tests.conftest import full_deck

ROOT = Path(__file__).resolve().parents[2]
DEFAULT_TEMPLATE_PATH = ROOT / "templates" / "default.pptx"
DEFAULT_MANIFEST_PATH = ROOT / "templates" / "default.manifest.json"
ENTERPRISE_EVAL_DIR = ROOT / "examples" / "enterprise-eval"
SLIDE_WIDTH_EMU = 12_192_000
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"p": PML_NS, "a": A_NS}
CORE_PROPS_NS = {
    "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
    "dc": "http://purl.org/dc/elements/1.1/",
    "dcterms": "http://purl.org/dc/terms/",
}
APP_PROPS_NS = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
SENSITIVE_TEMPLATE_METADATA_TERMS = (
    "SB",
    "SoftBank",
    "softbank",
    "ソフトバンク",
    "向山",
    "裕一朗",
    "法人事業統括",
)


def _slide_texts(slide: object) -> list[str]:
    return [
        shape.text.strip()
        for shape in getattr(slide, "shapes", [])
        if hasattr(shape, "text") and shape.text.strip()
    ]


def _slot_idx(manifest: dict[str, object], layout: str, slot_path: str) -> int:
    return int(manifest["layouts"][layout]["slots"][slot_path]["placeholder"]["idx"])


def _slot_text(slide: object, manifest: dict[str, object], layout: str, slot_path: str) -> str:
    return slide.placeholders[_slot_idx(manifest, layout, slot_path)].text.strip()


def _slot_first_paragraph(slide: object, manifest: dict[str, object], layout: str, slot_path: str) -> object:
    placeholder = slide.placeholders[_slot_idx(manifest, layout, slot_path)]
    return placeholder.text_frame.paragraphs[0]


def _slot_font_size(slide: object, manifest: dict[str, object], layout: str, slot_path: str) -> float:
    paragraph = _slot_first_paragraph(slide, manifest, layout, slot_path)
    assert paragraph.runs
    size = paragraph.runs[0].font.size
    assert size is not None
    return float(size.pt)


def _shape_by_name(root: etree._Element, name: str) -> etree._Element:
    matches = root.xpath(f".//p:sp[p:nvSpPr/p:cNvPr/@name='{name}']", namespaces=NS)
    assert len(matches) == 1
    return matches[0]


def _layout_roots() -> list[etree._Element]:
    with ZipFile(DEFAULT_TEMPLATE_PATH) as pptx:
        return [
            etree.fromstring(pptx.read(name))
            for name in pptx.namelist()
            if name.startswith("ppt/slideLayouts/slideLayout") and name.endswith(".xml")
        ]


def _shape_names(root: etree._Element) -> list[str]:
    return [
        shape.get("name")
        for shape in root.xpath(".//p:cNvPr[@name]", namespaces=NS)
        if shape.get("name")
    ]


def _shapes_by_name(roots: list[etree._Element], name: str) -> list[etree._Element]:
    return [
        shape
        for root in roots
        for shape in root.xpath(f".//p:sp[p:nvSpPr/p:cNvPr/@name='{name}']", namespaces=NS)
    ]


def _shape_extents(shape: etree._Element) -> tuple[int, int]:
    ext = shape.find("p:spPr/a:xfrm/a:ext", NS)
    assert ext is not None
    return int(ext.get("cx", "0")), int(ext.get("cy", "0"))


def _shape_line_is_empty(shape: etree._Element) -> bool:
    return shape.find("p:spPr/a:ln/a:noFill", NS) is not None


def _load_yaml_deck(name: str) -> dict[str, object]:
    return yaml.safe_load((ENTERPRISE_EVAL_DIR / name).read_text(encoding="utf-8"))


def _render_yaml_deck(name: str) -> bytes:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))
    return render_pptx(template_bytes, manifest, _load_yaml_deck(name))


def _rendered_slide_roots(pptx_bytes: bytes) -> list[etree._Element]:
    with ZipFile(BytesIO(pptx_bytes)) as pptx:
        return [
            etree.fromstring(pptx.read(name))
            for name in pptx.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]


def _ppt_xml_payload(pptx_bytes: bytes) -> bytes:
    with ZipFile(BytesIO(pptx_bytes)) as pptx:
        return b"\n".join(
            pptx.read(name)
            for name in pptx.namelist()
            if name.startswith("ppt/") and name.endswith(".xml")
        )


def _assert_sanitized_metadata(package_bytes: bytes, *, allow_chart_embeddings: bool = False) -> None:
    rejected_prefixes = (
        "customXml/",
        "ppt/comments",
        "ppt/commentAuthors",
        "ppt/people",
        "ppt/notesSlides/",
        "ppt/notesMasters/",
        "ppt/vbaProject.bin",
        "ppt/printerSettings/",
    )
    if not allow_chart_embeddings:
        rejected_prefixes = (*rejected_prefixes, "ppt/embeddings/")

    with ZipFile(BytesIO(package_bytes)) as pptx:
        names = set(pptx.namelist())
        core = etree.fromstring(pptx.read("docProps/core.xml"))
        app = etree.fromstring(pptx.read("docProps/app.xml"))

        assert not names & {"docProps/thumbnail.jpeg", "docProps/custom.xml", "EncryptionInfo", "EncryptedPackage"}
        assert not any(
            name.startswith(prefix)
            for name in names
            for prefix in rejected_prefixes
        )

        for term in SENSITIVE_TEMPLATE_METADATA_TERMS:
            assert not any(term.encode("utf-8") in pptx.read(name) for name in names)
            assert not any(term.encode("utf-16le") in pptx.read(name) for name in names)

    assert core.xpath("string(/cp:coreProperties/dc:creator)", namespaces=CORE_PROPS_NS) == ""
    assert core.xpath("string(/cp:coreProperties/cp:lastModifiedBy)", namespaces=CORE_PROPS_NS) == ""
    assert core.xpath("string(/cp:coreProperties/dc:description)", namespaces=CORE_PROPS_NS) == ""
    assert core.xpath("string(/cp:coreProperties/cp:revision)", namespaces=CORE_PROPS_NS) == "1"
    assert not core.xpath("/cp:coreProperties/dcterms:created", namespaces=CORE_PROPS_NS)
    assert not core.xpath("/cp:coreProperties/dcterms:modified", namespaces=CORE_PROPS_NS)
    assert app.xpath("string(/ep:Properties/ep:Application)", namespaces=APP_PROPS_NS) == ""
    assert app.xpath("string(/ep:Properties/ep:AppVersion)", namespaces=APP_PROPS_NS) == ""
    assert app.xpath("string(/ep:Properties/ep:Company)", namespaces=APP_PROPS_NS) == ""
    assert app.xpath("string(/ep:Properties/ep:Manager)", namespaces=APP_PROPS_NS) == ""


def test_repo_default_template_loads_via_registry() -> None:
    registry = TemplateRegistry(str(ROOT / "templates"))
    warnings = registry.load()

    assert warnings == []
    entry = registry.get_default()
    assert entry is not None
    assert entry.name == "default"
    assert sorted(entry.supported_layouts) == sorted(LAYOUT_SPECS.keys())


def test_repo_default_template_manifest_is_valid() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    checked_in_manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))
    manifest = generate_manifest(template_bytes)
    report = validate_manifest(template_bytes, manifest)

    assert manifest == checked_in_manifest
    assert set(manifest["layouts"].keys()) == set(LAYOUT_SPECS.keys())
    assert report["valid"], report["issues"]


def test_repo_default_template_uses_strict_ai_selection_pane_names() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    inspection = inspect_template(template_bytes)
    manifest = generate_manifest(template_bytes)
    layouts_by_name = {layout["layout_name"]: layout for layout in inspection["layouts"]}

    manifest_shape_names = [
        slot["placeholder"]["shape_name"]
        for binding in manifest["layouts"].values()
        for slot in binding["slots"].values()
    ]

    assert manifest_shape_names
    assert all(name.startswith("AI_") for name in manifest_shape_names)
    assert all("slot__" not in name and "placeholder__" not in name for name in manifest_shape_names)

    for semantic, binding in manifest["layouts"].items():
        allowed_names = {
            name
            for slot in LAYOUT_SPECS[semantic].slots
            for name in slot_ai_placeholder_names(semantic, slot.path)
        }
        layout = layouts_by_name[binding["ppt_layout_name"]]
        ai_placeholder_names = {
            placeholder["shape_name"]
            for placeholder in layout["placeholders"]
            if str(placeholder["shape_name"]).startswith("AI_")
        }
        bound_names = {slot["placeholder"]["shape_name"] for slot in binding["slots"].values()}

        assert bound_names <= ai_placeholder_names
        assert ai_placeholder_names <= allowed_names


def test_repo_default_template_uses_borderless_cards_and_square_icon_regions() -> None:
    roots = _layout_roots()

    for name in (
        "deco_content_panel",
        "deco_left_panel",
        "deco_right_panel",
        "deco_body_panel",
        "deco_refs_panel",
        "deco_media_panel",
        "deco_caption_panel",
        "deco_h_card1_bg",
        "deco_h_card2_bg",
        "deco_h_card3_bg",
    ):
        matches = _shapes_by_name(roots, name)
        assert matches
        assert all(_shape_line_is_empty(shape) for shape in matches)

    for name in (
        "deco_left_badge",
        "deco_right_badge",
        "deco_h_card1_badge",
        "deco_h_card2_badge",
        "deco_h_card3_badge",
        "AI_LEFT_ICON",
        "AI_RIGHT_ICON",
        "AI_CARDS_1_ICON",
        "AI_CARDS_2_ICON",
        "AI_CARDS_3_ICON",
        "AI_ICON",
    ):
        matches = _shapes_by_name(roots, name)
        assert matches
        for shape in matches:
            width, height = _shape_extents(shape)
            assert abs(width - height) <= 1


def test_repo_default_template_has_no_watermark_or_background_visual_fields() -> None:
    for root in _layout_roots():
        names = _shape_names(root)
        assert not any(name.startswith("deco_watermark_") for name in names)
        assert "deco_title_rule" not in names
        assert "deco_title_line" not in names
        assert "deco_section_line" not in names

    cover_or_closing_roots = [
        root
        for root in _layout_roots()
        if root.find("p:cSld", NS).get("name") in {"cover_title", "closing_end"}
    ]
    assert {root.find("p:cSld", NS).get("name") for root in cover_or_closing_roots} == {"cover_title", "closing_end"}
    for root in cover_or_closing_roots:
        assert "deco_visual_field" not in _shape_names(root)


def test_repo_default_cover_title_box_is_wide_enough_for_japanese_heading() -> None:
    cover_roots = [
        root
        for root in _layout_roots()
        if root.find("p:cSld", NS).get("name") == "cover_title"
    ]
    assert len(cover_roots) == 1
    root = cover_roots[0]
    title_shape = _shape_by_name(root, "AI_TITLE")

    title_width, _title_height = _shape_extents(title_shape)

    title_width_ratio = title_width / SLIDE_WIDTH_EMU

    assert title_width_ratio >= 0.575


def test_repo_default_template_clears_unspecified_optional_text_and_list_placeholders() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))
    deck = {
        "version": 1,
        "meta": {"title": "Optional Placeholder Cleanup"},
        "slides": [
            {"layout": "cover_title", "title": "Cover", "subtitle": "Subtitle"},
            {"layout": "list_basic", "title": "List", "items": ["One", "Two"]},
            {"layout": "closing_end", "title": "Thanks"},
        ],
    }

    output = render_pptx(template_bytes, manifest, deck)
    prs = Presentation(BytesIO(output))
    all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert "Text Placeholder" not in all_text
    assert "テキストを入力" not in all_text
    assert _slot_text(prs.slides[0], manifest, "cover_title", "date") == ""
    assert _slot_text(prs.slides[0], manifest, "cover_title", "organization") == ""
    assert _slot_text(prs.slides[0], manifest, "cover_title", "author") == ""
    assert _slot_text(prs.slides[1], manifest, "list_basic", "subtitle") == ""
    assert _slot_text(prs.slides[2], manifest, "closing_end", "subtitle") == ""
    assert _slot_text(prs.slides[2], manifest, "closing_end", "message") == ""
    assert _slot_text(prs.slides[2], manifest, "closing_end", "contact") == ""
    assert _slot_text(prs.slides[2], manifest, "closing_end", "cta") == ""


def test_repo_default_template_rendered_slide_shape_names_use_manifest_ai_names() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))

    output = render_pptx(template_bytes, manifest, full_deck())
    slide_shape_names = [
        name
        for root in _rendered_slide_roots(output)
        for name in _shape_names(root)
    ]

    assert not any("Text Placeholder" in name for name in slide_shape_names)


def test_repo_default_template_and_rendered_package_have_no_generic_placeholder_markers() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))
    rendered_bytes = render_pptx(template_bytes, manifest, full_deck())

    for package_bytes in (template_bytes, rendered_bytes):
        package_xml = _ppt_xml_payload(package_bytes)
        for marker in (
            b"Text Placeholder",
            "テキストを入力".encode(),
            b"deco_visual_field",
            b"deco_watermark",
            b"deco_title_rule",
            b"deco_title_line",
            b"deco_section_line",
        ):
            assert marker not in package_xml


def test_repo_default_template_and_rendered_package_have_sanitized_metadata() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))
    rendered_bytes = render_pptx(template_bytes, manifest, full_deck())

    _assert_sanitized_metadata(template_bytes)
    _assert_sanitized_metadata(rendered_bytes, allow_chart_embeddings=True)


@pytest.mark.parametrize(
    ("deck_name", "expected_language", "expected_font"),
    [
        ("default-enterprise-en.yaml", "en", FONT_FAMILY),
        ("default-enterprise-ja.yaml", "ja", CJK_FONT_FAMILY),
    ],
)
def test_repo_default_template_renders_enterprise_sample_decks(deck_name: str, expected_language: str, expected_font: str) -> None:
    deck = _load_yaml_deck(deck_name)
    output = _render_yaml_deck(deck_name)
    prs = Presentation(BytesIO(output))

    assert deck["meta"]["language"] == expected_language
    assert len(prs.slides) == len(deck["slides"])

    with ZipFile(BytesIO(output)) as pptx:
        slide_xml = b"\n".join(
            pptx.read(name)
            for name in pptx.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )

    assert expected_font.encode() in slide_xml
    if expected_language == "ja":
        assert b'<a:ea typeface="Yu Gothic"' in slide_xml
        assert b"rFonts" not in slide_xml
    else:
        assert CJK_FONT_FAMILY.encode() not in slide_xml


def test_repo_default_template_adjusts_japanese_cover_title_without_over_shrinking() -> None:
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))
    ja_prs = Presentation(BytesIO(_render_yaml_deck("default-enterprise-ja.yaml")))
    en_prs = Presentation(BytesIO(_render_yaml_deck("default-enterprise-en.yaml")))

    ja_cover_size = _slot_font_size(ja_prs.slides[0], manifest, "cover_title", "title")
    en_cover_size = _slot_font_size(en_prs.slides[0], manifest, "cover_title", "title")

    assert 36.0 <= ja_cover_size < 40.0
    assert en_cover_size == 40.0


@pytest.mark.parametrize("deck_name", ["default-enterprise-en.yaml", "default-enterprise-ja.yaml"])
def test_repo_default_enterprise_samples_stay_within_text_capacity_budget(deck_name: str) -> None:
    output = _render_yaml_deck(deck_name)
    prs = Presentation(BytesIO(output))
    issues: list[dict[str, object]] = []

    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame.text.strip():
                continue
            paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
            if len(paragraphs) > 7:
                issues.append({"slide": slide_index, "shape": shape.name, "paragraphs": len(paragraphs)})
            for paragraph in paragraphs:
                width = _text_display_width(paragraph.text)
                font_size = paragraph.runs[0].font.size.pt if paragraph.runs and paragraph.runs[0].font.size else None
                if width > 96:
                    issues.append({"slide": slide_index, "shape": shape.name, "width": width, "text": paragraph.text})
                if font_size is not None and font_size < 8.0:
                    issues.append({"slide": slide_index, "shape": shape.name, "font_size": font_size, "text": paragraph.text})

    assert issues == []


def test_repo_default_template_renders_full_deck() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))

    output = render_pptx(template_bytes, manifest, full_deck())
    prs = Presentation(BytesIO(output))

    assert len(prs.slides) == len(full_deck()["slides"])
    texts = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "Fixture Deck" not in texts
    assert "Cover" in texts


def test_repo_default_template_special_layouts_render_structured_content() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))

    output = render_pptx(template_bytes, manifest, full_deck())
    prs = Presentation(BytesIO(output))

    list_texts = _slide_texts(prs.slides[3])
    list_combined = "\n".join(list_texts)
    assert "A" in list_combined
    assert "B" in list_combined
    assert "1" not in list_texts

    comparison_slide = prs.slides[5]
    assert _slot_text(comparison_slide, manifest, "comparison_2col", "left.title") == "Old"
    assert _slot_text(comparison_slide, manifest, "comparison_2col", "left.description") == "Before"
    assert _slot_text(comparison_slide, manifest, "comparison_2col", "right.title") == "New"
    assert _slot_text(comparison_slide, manifest, "comparison_2col", "right.description") == "After"

    cards_v = prs.slides[6]
    assert _slot_text(cards_v, manifest, "three_cards_vertical", "cards[0].title") == "A"
    assert _slot_first_paragraph(cards_v, manifest, "three_cards_vertical", "cards[0].title").alignment == PP_ALIGN.CENTER
    assert _slot_text(cards_v, manifest, "three_cards_vertical", "cards[0].description") == "AA"
    assert _slot_text(cards_v, manifest, "three_cards_vertical", "cards[1].title") == "B"
    assert _slot_first_paragraph(cards_v, manifest, "three_cards_vertical", "cards[1].title").alignment == PP_ALIGN.CENTER
    assert _slot_text(cards_v, manifest, "three_cards_vertical", "cards[1].description") == "BB"
    assert _slot_text(cards_v, manifest, "three_cards_vertical", "cards[2].title") == "C"
    assert _slot_first_paragraph(cards_v, manifest, "three_cards_vertical", "cards[2].title").alignment == PP_ALIGN.CENTER
    assert _slot_text(cards_v, manifest, "three_cards_vertical", "cards[2].description") == "CC"

    table_slide = prs.slides[4]
    assert any(
        getattr(shape, "is_placeholder", False)
        and shape.placeholder_format.idx == _slot_idx(manifest, "table_basic", "table")
        for shape in table_slide.shapes
    )

    chart_slide = prs.slides[8]
    assert any(
        getattr(shape, "is_placeholder", False)
        and shape.placeholder_format.idx == _slot_idx(manifest, "chart_basic", "chart")
        for shape in chart_slide.shapes
    )

    image_slide = prs.slides[9]
    assert _slot_text(image_slide, manifest, "image_caption", "title") == "Icon"
    assert any(
        getattr(shape, "is_placeholder", False)
        and shape.placeholder_format.idx == _slot_idx(manifest, "image_caption", "icon")
        for shape in image_slide.shapes
    )

    appendix_slide = prs.slides[10]
    assert _slot_text(appendix_slide, manifest, "appendix_backup", "title") == "Appendix"
    assert _slot_text(appendix_slide, manifest, "appendix_backup", "items") == "Ref"


def test_repo_default_template_styles_charts_monochrome() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    manifest = json.loads(DEFAULT_MANIFEST_PATH.read_text(encoding="utf-8"))
    deck = {
        "version": 1,
        "meta": {"title": "Chart Styling"},
        "slides": [
            {
                "layout": "chart_basic",
                "title": "Chart",
                "chart": {
                    "kind": "column",
                    "categories": ["Q1", "Q2"],
                    "series": [
                        {"name": "Primary", "values": [10, 12]},
                        {"name": "Secondary", "values": [8, 9]},
                        {"name": "Tertiary", "values": [5, 7]},
                    ],
                },
            }
        ],
    }

    output = render_pptx(template_bytes, manifest, deck)
    prs = Presentation(BytesIO(output))
    chart = next(shape.chart for shape in prs.slides[0].shapes if getattr(shape, "has_chart", False))

    expected_colors = [RGBColor(0x19, 0x19, 0x19), RGBColor(0x24, 0x24, 0x24), RGBColor(0x74, 0x74, 0x74)]
    assert [series.format.fill.fore_color.rgb for series in chart.series] == expected_colors
    assert [series.format.line.color.rgb for series in chart.series] == expected_colors
    assert chart.category_axis.tick_labels.font.size == Pt(9)
    assert chart.value_axis.tick_labels.font.size == Pt(9)
    assert chart.category_axis.tick_labels.font.color.rgb == RGBColor(0x66, 0x66, 0x66)
    assert chart.value_axis.major_gridlines.format.line.color.rgb == RGBColor(0xD9, 0xD9, 0xD6)
