from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation

from pptx_yaml_engine import cli

ROOT = Path(__file__).resolve().parents[2]


def test_cli_generates_pptx_from_repo_default_assets_and_sample_yaml(tmp_path: Path) -> None:
    output_path = tmp_path / "sample-output.pptx"

    exit_code = cli.main(
        [
            "--template",
            str(ROOT / "templates" / "default.pptx"),
            "--manifest",
            str(ROOT / "templates" / "default.manifest.json"),
            "--deck",
            str(ROOT / "examples" / "review" / "sample-deck.yaml"),
            "--output",
            str(output_path),
        ]
    )

    assert exit_code == 0
    prs = Presentation(BytesIO(output_path.read_bytes()))
    assert len(prs.slides) == 11
    texts = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "Cover" in texts
    assert "Appendix" in texts
