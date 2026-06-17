from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation

from pptx_yaml_engine.mapper.service import generate_manifest, inspect_template, validate_manifest
from pptx_yaml_engine.output.service import render_pptx
from pptx_yaml_engine.server.template_registry import TemplateRegistry
from tests.conftest import full_deck, make_potx_bytes

ROOT = Path(__file__).resolve().parents[2]
DEFAULT_TEMPLATE_PATH = ROOT / "templates" / "default.pptx"


def _named_fixture_template_bytes() -> bytes:
    prs = Presentation(BytesIO(DEFAULT_TEMPLATE_PATH.read_bytes()))
    prs.core_properties.subject = "Named template registry fixture"
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def test_generate_manifest_and_render_all_layouts_with_default_template() -> None:
    template_bytes = DEFAULT_TEMPLATE_PATH.read_bytes()
    template_manifest = generate_manifest(template_bytes)
    inspection = inspect_template(template_bytes)
    assert inspection["template_fingerprint"] == template_manifest["template_fingerprint"]
    manifest_report = validate_manifest(template_bytes, template_manifest)
    assert manifest_report["valid"], manifest_report["issues"]

    output = render_pptx(template_bytes, template_manifest, full_deck())
    prs = Presentation(BytesIO(output))
    assert len(prs.slides) == len(full_deck()["slides"])
    texts = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "Cover" in texts
    assert "Appendix" in texts


def test_potx_template_bytes_are_accepted() -> None:
    inspection = inspect_template(make_potx_bytes())
    assert inspection["layouts"]


def test_named_custom_template_loads_and_renders_via_registry(tmp_path: Path) -> None:
    template_dir = tmp_path / "templates"
    template_dir.mkdir()
    custom_bytes = _named_fixture_template_bytes()
    (template_dir / "executive.pptx").write_bytes(custom_bytes)

    registry = TemplateRegistry(str(template_dir))
    warnings = registry.load()

    assert warnings == []
    entry = registry.resolve("Executive")
    assert entry is not None
    assert entry.name == "executive"
    assert entry.manifest["template_fingerprint"] == inspect_template(custom_bytes)["template_fingerprint"]
    assert validate_manifest(entry.template_bytes, entry.manifest)["valid"]

    output = render_pptx(entry.template_bytes, entry.manifest, full_deck())
    prs = Presentation(BytesIO(output))
    assert len(prs.slides) == len(full_deck()["slides"])
    assert prs.core_properties.subject == "Named template registry fixture"
    texts = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "Cover" in texts
    assert "Appendix" in texts
