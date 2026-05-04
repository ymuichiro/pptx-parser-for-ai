from __future__ import annotations

from io import BytesIO
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pptx import Presentation

from pptx_yaml_engine.layouts import LAYOUT_SPECS
from pptx_yaml_engine.mapper.service import finalize_manifest, inspect_template, propose_mapping
from pptx_yaml_engine.utils.template_bytes import POTX_MAIN_CONTENT_TYPE, PPTX_MAIN_CONTENT_TYPE


def make_template_bytes() -> bytes:
    prs = Presentation()
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def make_potx_bytes() -> bytes:
    source = BytesIO(make_template_bytes())
    output = BytesIO()
    with ZipFile(source, "r") as zin, ZipFile(output, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(PPTX_MAIN_CONTENT_TYPE, POTX_MAIN_CONTENT_TYPE)
            zout.writestr(item, data)
    return output.getvalue()


def _placeholder_idx(layout: dict[str, Any], needle: str | None = None) -> int:
    placeholders = layout["placeholders"]
    if needle is not None:
        for placeholder in placeholders:
            if needle in str(placeholder.get("placeholder_type", "")).upper():
                return int(placeholder["placeholder_idx"])
    return int(placeholders[0]["placeholder_idx"])


def _content_idx(layout: dict[str, Any]) -> int:
    title_idx = _placeholder_idx(layout, "TITLE")
    for placeholder in layout["placeholders"]:
        idx = int(placeholder["placeholder_idx"])
        if idx != title_idx:
            return idx
    return title_idx


def _binding(idx: int, kind: str = "text") -> dict[str, Any]:
    return {"idx": idx, "type": "BODY", "kind": kind}


def full_deck() -> dict[str, Any]:
    icon = {"pack": "heroicons", "name": "cloud-arrow-up", "variant": "outline"}
    return {
        "version": 1,
        "meta": {"title": "Fixture Deck", "language": "en"},
        "slides": [
            {"layout": "cover_title", "title": "Cover", "subtitle": "Subtitle"},
            {"layout": "section_divider", "title": "Section"},
            {"layout": "agenda", "title": "Agenda", "items": ["One", "Two"]},
            {"layout": "list_basic", "title": "List", "items": [{"text": "A", "level": 0}, {"text": "B", "level": 1}]},
            {"layout": "table_basic", "title": "Table", "table": {"headers": ["A", "B"], "rows": [[1, 2]]}},
            {
                "layout": "comparison_2col",
                "title": "Compare",
                "left": {"title": "Old", "description": "Before", "icon": icon},
                "right": {"title": "New", "description": "After", "icon": icon},
            },
            {
                "layout": "three_cards_vertical",
                "title": "Cards V",
                "cards": [
                    {"title": "A", "description": "AA", "icon": icon},
                    {"title": "B", "description": "BB", "icon": icon},
                    {"title": "C", "description": "CC", "icon": icon},
                ],
            },
            {
                "layout": "three_cards_horizontal",
                "title": "Cards H",
                "cards": [
                    {"title": "A", "description": "AA", "icon": icon},
                    {"title": "B", "description": "BB", "icon": icon},
                    {"title": "C", "description": "CC", "icon": icon},
                ],
            },
            {"layout": "timeline", "title": "Timeline", "events": [{"label": "Q1", "title": "Plan", "description": "Start"}]},
            {"layout": "closing_end", "title": "Thanks"},
            {"layout": "kpi_big_number", "title": "KPI", "metric": {"value": "42", "label": "Customers"}},
            {"layout": "chart_basic", "title": "Chart", "chart": {"kind": "column", "categories": ["Q1", "Q2"], "series": [{"name": "Sales", "values": [10, 12]}]}},
            {"layout": "image_caption", "title": "Icon", "icon": icon},
            {"layout": "appendix_backup", "title": "Appendix", "items": ["Ref"]},
            {"layout": "eol_notice", "title": "EOL", "product_name": "Legacy API", "actions": ["Migrate"]},
        ],
    }


@pytest.fixture()
def template_bytes() -> bytes:
    return make_template_bytes()


@pytest.fixture()
def template_manifest(template_bytes: bytes) -> dict[str, Any]:
    inspection = inspect_template(template_bytes)
    title_slide = inspection["layouts"][0]
    content_slide = next(layout for layout in inspection["layouts"] if len(layout["placeholders"]) >= 2)
    title_idx = _placeholder_idx(content_slide, "TITLE")
    content_idx = _content_idx(content_slide)
    proposal = propose_mapping(
        inspection,
        {"aliases": {semantic: [content_slide["layout_name"]] for semantic in LAYOUT_SPECS}},
    )
    overrides: dict[str, Any] = {"layouts": {}}
    for semantic, spec in LAYOUT_SPECS.items():
        layout = title_slide if semantic == "cover_title" else content_slide
        slide_title_idx = _placeholder_idx(layout, "TITLE")
        slide_content_idx = _content_idx(layout)
        slots: dict[str, Any] = {}
        for slot in spec.slots:
            if not slot.required and slot.path not in {"subtitle", "items", "table", "chart", "icon", "product_name", "actions", "metric.label"}:
                continue
            idx = slide_title_idx if slot.path == "title" else slide_content_idx
            slots[slot.path] = _binding(idx, slot.kind)
        if semantic == "cover_title":
            slots["title"] = _binding(_placeholder_idx(title_slide, "TITLE"))
            slots["subtitle"] = _binding(_content_idx(title_slide))
        overrides["layouts"][semantic] = {
            "ppt_layout_name": layout["layout_name"],
            "slots": slots,
        }
    manifest = finalize_manifest(inspection, proposal, overrides)
    assert title_idx >= 0
    assert content_idx >= 0
    return manifest
