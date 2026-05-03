from __future__ import annotations

import json
from io import BytesIO
from pathlib import Path
from typing import Any

import pytest
from pptx import Presentation

from pptx_yaml_engine import cli


def _write_text(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def test_cli_generates_pptx_from_template_manifest_and_yaml(
    tmp_path: Path, template_bytes: bytes, template_manifest: dict[str, Any]
) -> None:
    template_path = tmp_path / "template.pptx"
    manifest_path = tmp_path / "manifest.json"
    deck_path = tmp_path / "deck.yaml"
    output_path = tmp_path / "result.pptx"
    template_path.write_bytes(template_bytes)
    manifest_path.write_text(json.dumps(template_manifest), encoding="utf-8")
    _write_text(
        deck_path,
        "version: 1\nslides:\n  - layout: cover_title\n    title: Cover\n    subtitle: Subtitle\n",
    )

    exit_code = cli.main(
        [
            "--template",
            str(template_path),
            "--manifest",
            str(manifest_path),
            "--deck",
            str(deck_path),
            "--output",
            str(output_path),
        ]
    )

    assert exit_code == 0
    prs = Presentation(BytesIO(output_path.read_bytes()))
    assert len(prs.slides) == 1
    texts = [shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text")]
    assert "Cover" in texts
    assert "Subtitle" in texts


def test_cli_rejects_invalid_yaml(
    tmp_path: Path,
    template_bytes: bytes,
    template_manifest: dict[str, Any],
    capsys: pytest.CaptureFixture[str],
) -> None:
    template_path = tmp_path / "template.pptx"
    manifest_path = tmp_path / "manifest.json"
    deck_path = tmp_path / "broken.yaml"
    output_path = tmp_path / "result.pptx"
    template_path.write_bytes(template_bytes)
    manifest_path.write_text(json.dumps(template_manifest), encoding="utf-8")
    _write_text(deck_path, "slides: [\n")

    exit_code = cli.main(
        [
            "--template",
            str(template_path),
            "--manifest",
            str(manifest_path),
            "--deck",
            str(deck_path),
            "--output",
            str(output_path),
        ]
    )

    captured = capsys.readouterr()
    assert exit_code == 1
    assert "DECK_YAML_INVALID" in captured.err
    assert not output_path.exists()


def test_cli_rejects_manifest_template_fingerprint_mismatch(
    tmp_path: Path,
    template_bytes: bytes,
    template_manifest: dict[str, Any],
    capsys: pytest.CaptureFixture[str],
) -> None:
    template_path = tmp_path / "template.pptx"
    manifest_path = tmp_path / "manifest.json"
    deck_path = tmp_path / "deck.yaml"
    output_path = tmp_path / "result.pptx"
    broken_manifest = dict(template_manifest)
    broken_manifest["template_fingerprint"] = "sha256:not-the-template"
    template_path.write_bytes(template_bytes)
    manifest_path.write_text(json.dumps(broken_manifest), encoding="utf-8")
    _write_text(
        deck_path,
        "version: 1\nslides:\n  - layout: cover_title\n    title: Cover\n    subtitle: Subtitle\n",
    )

    exit_code = cli.main(
        [
            "--template",
            str(template_path),
            "--manifest",
            str(manifest_path),
            "--deck",
            str(deck_path),
            "--output",
            str(output_path),
        ]
    )

    captured = capsys.readouterr()
    assert exit_code == 1
    assert "MANIFEST_VALIDATION_FAILED" in captured.err
    assert "TEMPLATE_FINGERPRINT_MISMATCH" in captured.err
    assert not output_path.exists()


def test_cli_rejects_missing_output_directory(
    tmp_path: Path,
    template_bytes: bytes,
    template_manifest: dict[str, Any],
    capsys: pytest.CaptureFixture[str],
) -> None:
    template_path = tmp_path / "template.pptx"
    manifest_path = tmp_path / "manifest.json"
    deck_path = tmp_path / "deck.yaml"
    output_path = tmp_path / "missing" / "result.pptx"
    template_path.write_bytes(template_bytes)
    manifest_path.write_text(json.dumps(template_manifest), encoding="utf-8")
    _write_text(
        deck_path,
        "version: 1\nslides:\n  - layout: cover_title\n    title: Cover\n    subtitle: Subtitle\n",
    )

    exit_code = cli.main(
        [
            "--template",
            str(template_path),
            "--manifest",
            str(manifest_path),
            "--deck",
            str(deck_path),
            "--output",
            str(output_path),
        ]
    )

    captured = capsys.readouterr()
    assert exit_code == 1
    assert "OUTPUT_PATH_INVALID" in captured.err
